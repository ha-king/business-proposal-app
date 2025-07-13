from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation
from pptx.util import Inches
import io
from datetime import datetime

def generate_pdf(proposal_content: str, market_content: str) -> bytes:
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    
    # Title
    title = Paragraph("Business Proposal & Market Analysis", styles['Title'])
    story.append(title)
    story.append(Spacer(1, 12))
    
    # Date
    date = Paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y')}", styles['Normal'])
    story.append(date)
    story.append(Spacer(1, 24))
    
    # Proposal Section
    proposal_title = Paragraph("Business Proposal", styles['Heading1'])
    story.append(proposal_title)
    story.append(Spacer(1, 12))
    
    proposal_text = Paragraph(proposal_content.replace('\n', '<br/>'), styles['Normal'])
    story.append(proposal_text)
    story.append(Spacer(1, 24))
    
    # Market Analysis Section
    market_title = Paragraph("Market Analysis", styles['Heading1'])
    story.append(market_title)
    story.append(Spacer(1, 12))
    
    market_text = Paragraph(market_content.replace('\n', '<br/>'), styles['Normal'])
    story.append(market_text)
    
    doc.build(story)
    buffer.seek(0)
    return buffer.getvalue()

def generate_pptx(proposal_content: str, market_content: str) -> bytes:
    prs = Presentation()
    
    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Business Proposal & Market Analysis"
    subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
    
    # Proposal slides
    sections = proposal_content.split('\n\n')
    for section in sections[:5]:  # Limit to 5 sections
        if section.strip():
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            lines = section.strip().split('\n')
            title.text = lines[0] if lines else "Section"
            content.text = '\n'.join(lines[1:]) if len(lines) > 1 else section
    
    # Market analysis slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Market Analysis"
    content.text = market_content[:500] + "..." if len(market_content) > 500 else market_content
    
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()